import uuid
from pptx import Presentation as pptx_lib
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx_rag_quizzer.presentation_model import Image, Text, Slide, Presentation, Type
from typing import List, Union

def parse_powerpoint(file_object, file_name):
    """
    Parses an in-memory PowerPoint file to extract text and images in order.

    The function extracts text from speaker notes and shapes, and image data
    from picture shapes, maintaining the slide order.

    Args:
        file_object (io.BytesIO): An in-memory byte stream of the .pptx file.

    Returns:
        Presentation: A Presentation object containing the slides.
    """
    prs = pptx_lib(file_object)

    PRESENTATION = Presentation(
        id=str(uuid.uuid4()),
        name=file_name,
        slides=[],
    )


    for slide_idx, slide in enumerate(prs.slides):
        slide_items: List[Union[Image, Text]] = []

        order_number = 0

        # Extract from speaker notes first
        if (
            slide.has_notes_slide
            and slide.notes_slide.notes_text_frame
            and slide.notes_slide.notes_text_frame.text
        ):
            text = Text(
                id=str(uuid.uuid4()),
                content=slide.notes_slide.notes_text_frame.text,
                slide_number=slide_idx + 1,
                type=Type.text,
                order_number=order_number,
            )
            slide_items.append(text)

        # Sort shapes by position (top-to-bottom, left-to-right) for reading order
        shapes = sorted(
            slide.shapes,
            key=lambda x: (
                (x.top, x.left) if hasattr(x, "top") and hasattr(x, "left") else (0, 0)
            ),
        )

        # Extract from shapes on the slide
        for shape in shapes:
            if shape.has_text_frame and shape.text_frame.text:
                text = Text(
                    id=str(uuid.uuid4()),
                    content=shape.text_frame.text,
                    slide_number=slide_idx + 1,
                    type=Type.text,
                    order_number=order_number,
                )
                slide_items.append(text)
                order_number += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_bytes = shape.image.blob
                image_ext = shape.image.ext
                image = Image(
                    id=str(uuid.uuid4()),
                    content='none',
                    extension=image_ext,
                    image_bytes=image_bytes,
                    slide_number=slide_idx + 1,
                    type=Type.image,
                    order_number=order_number,
                )
                slide_items.append(image)
                order_number += 1

        PRESENTATION.slides.append(Slide(
            id=str(uuid.uuid4()),
            slide_number=slide_idx + 1,
            items=slide_items,
        ))

    return PRESENTATION